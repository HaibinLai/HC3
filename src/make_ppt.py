#!/usr/bin/env python3
"""Generate PPT for AIGC Detection final project presentation (10 min)."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path
import os

BASE = Path(__file__).resolve().parent.parent
FIG = BASE / "figures"
OUT = BASE / "reports"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)
BG_MID  = RGBColor(0x16, 0x21, 0x3E)
ACCENT  = RGBColor(0x00, 0xD2, 0xFF)
ACCENT2 = RGBColor(0xFF, 0x6B, 0x6B)
ACCENT3 = RGBColor(0x4E, 0xCB, 0x71)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GRAY    = RGBColor(0xAA, 0xAA, 0xAA)
LIGHT   = RGBColor(0xE0, 0xE0, 0xE0)

def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=WHITE, spacing=1.2):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(font_size * spacing * 0.5)
        p.level = 0
    return txBox

def add_image_safe(slide, img_path, left, top, width=None, height=None):
    p = Path(img_path)
    if not p.exists():
        add_text_box(slide, left, top, 4, 1, f"[Image: {p.name}]", 12, GRAY)
        return
    if width and height:
        slide.shapes.add_picture(str(p), Inches(left), Inches(top), Inches(width), Inches(height))
    elif width:
        slide.shapes.add_picture(str(p), Inches(left), Inches(top), Inches(width))
    else:
        slide.shapes.add_picture(str(p), Inches(left), Inches(top))

def add_accent_line(slide, left, top, width, color=ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

# ========== Slide 1: Title ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide)
add_text_box(slide, 1, 1.5, 11, 1.5,
    "AI 生成文本的特征挖掘与检测", 44, ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
add_accent_line(slide, 3, 3.1, 7)
add_text_box(slide, 1, 3.5, 11, 1,
    "Feature Mining and Detection of AI-Generated Text", 24, LIGHT, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1, 4.8, 11, 0.8,
    "数据挖掘课程期末项目 (CS306) — 方向一", 20, GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1, 5.8, 11, 0.8,
    "赖海斌  |  2026 年 6 月", 18, GRAY, alignment=PP_ALIGN.CENTER)

# ========== Slide 2: Problem ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, 0.8, 0.4, 10, 0.8, "问题定义", 36, ACCENT, bold=True)
add_accent_line(slide, 0.8, 1.1, 4)
add_bullet_list(slide, 0.8, 1.5, 5.5, 5, [
    "核心任务：区分人类文本 vs AI 生成文本 (二分类)",
    "",
    "挑战 1: 多模型 — 11+ 种 AI 生成器，特征各异",
    "挑战 2: 多领域 — 新闻、诗歌、学术摘要等差异大",
    "挑战 3: 对抗攻击 — 释义、同义词替换等规避手段",
    "挑战 4: 可解释性 — 不能只给判定，要说明为什么",
    "",
    "我们的目标：",
    "  用可解释的统计特征 + 轻量机器学习",
    "  替代不可解释的深度学习黑盒",
], font_size=18)
add_image_safe(slide, FIG / "class_distribution.png", 7.5, 1.5, width=5)

# ========== Slide 3: Datasets ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, 0.8, 0.4, 10, 0.8, "数据集总览", 36, ACCENT, bold=True)
add_accent_line(slide, 0.8, 1.1, 4)

# Table-like layout
datasets = [
    ("RAID (主实验)", "ACL 2024", "467K", "11 生成器, 8 领域, 11 攻击", ACCENT),
    ("HC3", "对比实验", "85K", "ChatGPT 单模型, 4 QA 领域", LIGHT),
    ("SemEval 2024", "辅助验证", "154K", "多模型多领域", LIGHT),
    ("TuringBench", "辅助验证", "441K", "19 模型 (GPT-1~3 时代)", LIGHT),
    ("Pile", "辅助验证", "1.39M", "GPT-2 生成, 混合领域", LIGHT),
]
for i, (name, role, size, desc, clr) in enumerate(datasets):
    y = 1.8 + i * 1.0
    add_text_box(slide, 1.0, y, 2.5, 0.5, name, 20, clr, bold=True)
    add_text_box(slide, 3.6, y, 1.5, 0.5, role, 16, GRAY)
    add_text_box(slide, 5.2, y, 1.2, 0.5, size, 16, WHITE)
    add_text_box(slide, 6.5, y, 6, 0.5, desc, 16, WHITE)

add_text_box(slide, 0.8, 6.5, 11, 0.5,
    "RAID 是目前最全面的 AI 文本检测基准，覆盖从 GPT-2 到 GPT-4 的全代际生成器",
    14, GRAY, alignment=PP_ALIGN.CENTER)

# ========== Slide 4: Method Overview ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, 0.8, 0.4, 12, 0.8, "方法概览：三层 120 维特征体系 + XGBoost", 32, ACCENT, bold=True)
add_accent_line(slide, 0.8, 1.1, 8)

ACCENT_ORANGE = RGBColor(0xFF, 0xA5, 0x00)

# Column 1: Statistical
add_text_box(slide, 0.3, 1.4, 3.8, 0.5, "Layer 1: 纯统计特征 (~36维)", 18, ACCENT2, bold=True)
add_bullet_list(slide, 0.3, 1.9, 3.8, 3.2, [
    "basic_counts (4): 词/字符/句/段数",
    "averages (4): 平均词长/句长/段长",
    "variability (2): 词长/句长标准差",
    "lexical_richness (7): TTR, Hapax...",
    "punctuation (11): 标点 + 大写/数字比",
    "readability (7): Flesch, SMOG...",
    "structure (1): 短句比例",
    "",
    "不依赖任何神经网络, 纯公式计算",
], font_size=12, color=LIGHT)

# Column 2: Model-derived
add_text_box(slide, 4.5, 1.4, 4, 0.5, "Layer 2: 模型衍生特征 (52维)", 18, ACCENT_ORANGE, bold=True)
add_bullet_list(slide, 4.5, 1.9, 4, 3.2, [
    "BERT embedding → PCA (50维)",
    "  将文本映射到语义空间",
    "  AI 文本语义分布更集中",
    "",
    "GPT-2 perplexity (2维)",
    "  文本级别的困惑度汇总",
    "  HC3 最强特征 (AUC 0.99)",
    "  RAID 完全失效 (AUC 0.49)",
    "  → 模型依赖性是其本质属性",
], font_size=12, color=LIGHT)

# Column 3: Token-level
add_text_box(slide, 8.8, 1.4, 4.2, 0.5, "Layer 3: Token 概率特征 (30维)", 18, ACCENT3, bold=True)
add_bullet_list(slide, 8.8, 1.9, 4.2, 3.2, [
    "观察者: Mistral-7B-Instruct",
    "",
    "Log-prob (10): mean, std, min,",
    "  max, median, q10, q90, skew...",
    "Rank (8): mean, std, median,",
    "  top-1/5/10/100 fraction",
    "Entropy (6): mean, std, min,",
    "  max, median, skew",
    "",
    "逐 token 粒度, 看 AI 的\"选词指纹\"",
], font_size=12, color=LIGHT)

# Bottom: pipeline arrow
add_text_box(slide, 0.5, 5.3, 12, 0.5,
    "纯统计 (36维)  +  模型衍生 (52维)  +  Token级 (30维)  =  118维 (去重后)",
    18, WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 2, 5.8, 9, 0.6,
    "→  Auto-Filter (丢噪声组)  →  XGBoost (500 trees, depth 8)  →  SHAP 可解释分析",
    18, WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, 0.8, 6.5, 11, 0.5,
    "三层特征从不同粒度和视角刻画文本: 表面统计 → 语义嵌入 → Token 概率",
    16, ACCENT, alignment=PP_ALIGN.CENTER)

# ========== Slide 4b: Lexical Richness & Readability Deep Dive ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, 0.8, 0.4, 12, 0.8, "手工特征详解：词汇丰富度 & 可读性", 32, ACCENT, bold=True)
add_accent_line(slide, 0.8, 1.1, 8)

# Left column: Lexical Richness
add_text_box(slide, 0.8, 1.4, 5.5, 0.6, "Lexical Richness (词汇丰富度) — 7 维", 20, ACCENT2, bold=True)
add_bullet_list(slide, 0.8, 2.0, 5.8, 4.5, [
    "Type-Token Ratio (TTR)",
    "  不重复词数 / 总词数, 衡量词汇多样性",
    "  AI 倾向反复使用\"furthermore, however, overall\"",
    "",
    "Hapax Legomena Ratio",
    "  只出现一次的词占比, 人类用词更\"独特\"",
    "",
    "Yule's K & Simpson's Diversity",
    "  基于词频分布的数学度量",
    "  Yule's K 越大 → 词汇越集中(越像 AI)",
    "",
    "Brunet's W",
    "  W = N^(V^-0.172), 对文本长度鲁棒的丰富度指标",
], font_size=14, color=LIGHT)

# Right column: Readability
add_text_box(slide, 7, 1.4, 5.5, 0.6, "Readability (可读性) — 7 维", 20, ACCENT3, bold=True)
add_bullet_list(slide, 7, 2.0, 5.8, 4.5, [
    "Flesch Reading Ease  (0-100)",
    "  206.835 - 1.015×(词/句) - 84.6×(音节/词)",
    "  AI 文本通常 50-65 分, 人类分布更广",
    "",
    "Flesch-Kincaid Grade Level",
    "  输出\"需要几年级才能读懂\", AI 偏好 9-12 级",
    "",
    "Gunning Fog Index",
    "  关注\"复杂词\"(≥3音节)的比例",
    "",
    "Coleman-Liau Index",
    "  基于字符数而非音节, 更适合计算机处理",
    "",
    "SMOG, ARI, Dale-Chall",
    "  不同侧面的可读性评估",
], font_size=14, color=LIGHT)

add_text_box(slide, 0.8, 6.5, 12, 0.5,
    "AI 倾向于生成\"中等难度\"的文本 (Flesch 50-65), 而人类文本的可读性分布更极端 (20-90+)",
    15, ACCENT, alignment=PP_ALIGN.CENTER)

# ========== Slide 4c: Structural & Punctuation Features ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, 0.8, 0.4, 12, 0.8, "模型衍生特征详解：语义嵌入 & 困惑度", 32, ACCENT, bold=True)
add_accent_line(slide, 0.8, 1.1, 8)

# Left: Structure + Punctuation
add_text_box(slide, 0.8, 1.4, 5.8, 0.6, "结构 & 标点 (纯统计特征) — 16 维", 20, ACCENT2, bold=True)
add_bullet_list(slide, 0.8, 2.0, 5.8, 2.5, [
    "words_per_paragraph  (RAID 最强单特征, AUC 0.89)",
    "  AI 平均每段 ~70 词, 人类 ~200 词",
    "  AI 喜欢\"分点陈述\", 段落短而规则",
    "",
    "sentence_length_std  (AUC 0.83)",
    "  AI 的句子长度标准差更小 → 节奏更\"机械\"",
    "  人类写作句长波动大 (短句+长句交替)",
    "",
    "标点使用率: 逗号、问号、感叹号、冒号...",
    "  AI 很少用感叹号, 问号使用更规范",
    "  人类更随意: \"!!!\" \"...\" \"—\"",
], font_size=14, color=LIGHT)

# Right: Embedding PCA + Perplexity
add_text_box(slide, 7, 1.4, 5.8, 0.6, "模型衍生特征 (Layer 2) — 52 维", 20, RGBColor(0xFF, 0xA5, 0x00), bold=True)
add_bullet_list(slide, 7, 2.0, 5.8, 2.5, [
    "BERT Sentence Embedding → PCA 50 维",
    "  将文本映射到 BERT 语义空间",
    "  PCA 降维后保留主要语义方向",
    "  AI 文本在语义空间中更\"集中\"",
    "",
    "GPT-2 Perplexity (2 维)",
    "  原始困惑度 + log 困惑度",
    "  HC3 上是最强特征 (AUC 0.99)",
    "  但 RAID 上完全失效 (AUC 0.49)",
    "  → Perplexity Paradox (后面详述)",
], font_size=14, color=LIGHT)

# Bottom: boxplot
add_image_safe(slide, FIG / "raid_single_feature_boxplot.png", 0.5, 4.8, width=12, height=2.5)

# ========== Slide 4d: Token Feature Case Study ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, 0.8, 0.4, 12, 0.8, "Token 概率特征：Case Study 热力图", 32, ACCENT, bold=True)
add_accent_line(slide, 0.8, 1.1, 8)

add_image_safe(slide, FIG / "raid_token_heatmap.png", 0.2, 1.2, width=7.5, height=5.5)

add_text_box(slide, 8, 1.4, 5, 0.5, "如何阅读热力图", 20, ACCENT, bold=True)
add_bullet_list(slide, 8, 2.0, 5, 4.5, [
    "每个格子 = 一个 token",
    "颜色 = Mistral-7B 给出的 log-prob",
    "  深绿 = 高概率 (\"意料之中\")",
    "  深红 = 低概率 (\"出乎意料\")",
    "",
    "人类文本 (上半部分):",
    "  红色格子多 → 用词\"出人意料\"",
    "  log-prob 分布更分散",
    "",
    "AI 文本 (下半部分):",
    "  几乎全绿 → token 选择\"可预测\"",
    "  AI 倾向选择高概率的 token",
    "",
    "→ 这就是 rank_top1_frac 能区分",
    "  AI vs 人类的直觉来源",
], font_size=14, color=LIGHT)

add_text_box(slide, 0.8, 6.8, 12, 0.5,
    "AI 生成的文本在 Mistral 眼中\"一路绿灯\" — 每个 token 都是高概率选择",
    15, ACCENT, alignment=PP_ALIGN.CENTER)

# ========== Slide 4e: Token Feature Distributions ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, 0.8, 0.4, 12, 0.8, "Token 特征的统计分布：Human vs AI", 32, ACCENT, bold=True)
add_accent_line(slide, 0.8, 1.1, 8)

add_image_safe(slide, FIG / "raid_token_distributions.png", 0.3, 1.2, width=7, height=5)

add_text_box(slide, 7.8, 1.4, 5, 0.5, "核心 Token 特征解读", 20, ACCENT, bold=True)
add_bullet_list(slide, 7.8, 2.0, 5, 5, [
    "rank_top1_frac (top-1 命中率)",
    "  AI: ~0.72  人类: ~0.58",
    "  AI 有 72% 的 token 是 Mistral 的首选",
    "",
    "lp_mean (平均 log-probability)",
    "  AI: ~ -1.8  人类: ~ -2.8",
    "  AI 文本整体\"更可预测\"",
    "",
    "ent_mean (平均 entropy)",
    "  AI: ~1.5  人类: ~2.2",
    "  AI 选词更\"确定\", 人类更\"随机\"",
    "",
    "这 3 个特征在 RAID 的 SHAP 分析中",
    "都位列 Top-15, 与手工特征互补",
], font_size=14, color=LIGHT)

add_text_box(slide, 0.8, 6.5, 12, 0.5,
    "Violin plot 清晰展示: AI 的 token 概率分布更\"确定\"(高 lp, 高 rank, 低 entropy)",
    15, ACCENT, alignment=PP_ALIGN.CENTER)

# ========== Slide 5: Main Results ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, 0.8, 0.4, 10, 0.8, "RAID 主实验结果", 36, ACCENT, bold=True)
add_accent_line(slide, 0.8, 1.1, 4)

add_image_safe(slide, FIG / "raid_final_comparison.png", 0.5, 1.5, width=6.5)

# Results text
add_text_box(slide, 7.5, 1.5, 5, 0.6, "关键数字", 24, ACCENT, bold=True)
results_items = [
    "XGBoost (120 组合):  AUC 0.9992",
    "XGBoost (90 手工):   AUC 0.9951",
    "XGBoost (30 Token):  AUC 0.9900",
    "Random Forest (90):  AUC 0.9923",
    "LR (90):             AUC 0.9653",
    "Fast-DetectGPT:      AUC 0.7815",
    "",
    "120 组合 vs 90 手工:",
    "  AUC +0.0041, Acc 99.15%",
    "",
    "完全可解释 + 超越零样本方法",
]
add_bullet_list(slide, 7.5, 2.3, 5, 4.5, results_items, font_size=16, color=LIGHT)

# ========== Slide 6: Single Feature AUC ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, 0.8, 0.4, 10, 0.8, "单特征判别力分析", 36, ACCENT, bold=True)
add_accent_line(slide, 0.8, 1.1, 4)
add_image_safe(slide, FIG / "raid_single_feature_auc.png", 0.3, 1.3, width=6.5)
add_image_safe(slide, FIG / "raid_single_feature_dist.png", 6.8, 1.3, width=6.2)

add_text_box(slide, 0.8, 6.6, 11, 0.5,
    "words_per_paragraph (AUC 0.89) 是最强单特征 — AI 倾向于生成规则的段落结构",
    15, ACCENT, alignment=PP_ALIGN.CENTER)

# ========== Slide 7: Feature Ablation ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, 0.8, 0.4, 10, 0.8, "特征组消融实验", 36, ACCENT, bold=True)
add_accent_line(slide, 0.8, 1.1, 4)
add_image_safe(slide, FIG / "raid_feature_ablation.png", 0.3, 1.3, width=6.5)

add_text_box(slide, 7.2, 1.5, 5.5, 0.6, "关键发现", 22, ACCENT, bold=True)
add_bullet_list(slide, 7.2, 2.2, 5.5, 4, [
    "basic_counts (4维) 单组即达 AUC 0.97",
    "  → 文本长度是最稳定的跨模型信号",
    "",
    "perplexity (2维) AUC 仅 0.49",
    "  → GPT-2 困惑度在多模型场景完全失效",
    "  → 但在 HC3 (ChatGPT) 上是最强特征",
    "",
    "embedding_pca (50维) AUC 0.94",
    "  → BERT 语义嵌入捕捉了细微风格差异",
], font_size=15, color=LIGHT)

# ========== Slide 8: SHAP ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, 0.8, 0.4, 10, 0.8, "SHAP 可解释性分析 (120 维)", 36, ACCENT, bold=True)
add_accent_line(slide, 0.8, 1.1, 4)
add_image_safe(slide, FIG / "raid_120_shap_summary.png", 0.3, 1.3, width=6)
add_image_safe(slide, FIG / "raid_shap_dependence.png", 6.5, 1.3, width=6.5)
add_text_box(slide, 0.8, 6.6, 11, 0.5,
    "SHAP 揭示两类特征交错出现在 Top-30 — XGBoost 同时利用两个视角做决策",
    15, ACCENT, alignment=PP_ALIGN.CENTER)

# ========== Slide 9: Perplexity Paradox ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, 0.8, 0.4, 10, 0.8, "Perplexity Paradox: 困惑度的两面性", 32, ACCENT2, bold=True)
add_accent_line(slide, 0.8, 1.1, 6)

add_text_box(slide, 0.8, 1.5, 5.5, 0.6, "HC3 (ChatGPT 单模型)", 22, ACCENT3, bold=True)
add_bullet_list(slide, 0.8, 2.1, 5.5, 2, [
    "GPT-2 Perplexity AUC = 0.9912 (最强特征!)",
    "ChatGPT 与 GPT-2 同属 GPT 家族",
    "共享 BPE tokenizer 和相似概率分布",
    "困惑度能精准区分\"GPT 系输出\"与人类",
], font_size=16, color=LIGHT)

add_text_box(slide, 0.8, 4.2, 5.5, 0.6, "RAID (11 个生成器)", 22, ACCENT2, bold=True)
add_bullet_list(slide, 0.8, 4.8, 5.5, 2, [
    "GPT-2 Perplexity AUC = 0.4920 (比随机还差!)",
    "Llama, MPT, Cohere 的概率分布与 GPT 完全不同",
    "困惑度在方向上甚至是反的",
    "→ 困惑度不是检测\"AI 生成\"，是检测\"GPT 系生成\"",
], font_size=16, color=LIGHT)

add_image_safe(slide, FIG / "domain_feature_comparison.png", 7, 1.5, width=5.8)

# ========== Slide 10: Feature Complementarity ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, 0.8, 0.4, 11, 0.8, "两类特征的互补机制", 36, ACCENT, bold=True)
add_accent_line(slide, 0.8, 1.1, 6)

# Table header
y = 1.5
for col, x, w in [("维度", 0.8, 2), ("手工特征", 3, 4.5), ("Token 特征", 7.8, 5)]:
    add_text_box(slide, x, y, w, 0.5, col, 18, ACCENT, bold=True)

rows = [
    ("观察角度", "\"文本长什么样\"", "\"文本在 LLM 眼中像什么\""),
    ("信号来源", "宏观结构 / 风格 / 词汇", "微观 token 概率分布"),
    ("物理类比", "用肉眼看一幅画", "用光谱仪分析颜料成分"),
    ("强项", "老 AI + 跨模型稳定信号", "同代新 AI (概率差异明显)"),
    ("弱项", "新 AI (表面统计趋近人类)", "老 AI (跨代概率失配)"),
]
for i, (dim, hc, tok) in enumerate(rows):
    y = 2.2 + i * 0.85
    add_text_box(slide, 0.8, y, 2, 0.5, dim, 16, ACCENT2 if i >= 3 else WHITE, bold=True)
    add_text_box(slide, 3, y, 4.5, 0.5, hc, 15, LIGHT)
    add_text_box(slide, 7.8, y, 5, 0.5, tok, 15, LIGHT)

add_text_box(slide, 0.8, 6.3, 11.5, 0.8,
    "互补效果: Cohere 检测率 90.4% → 94.5% (token 补上手工盲区), GPT-2 手工已 >99% (兜底老模型)",
    16, ACCENT3, alignment=PP_ALIGN.CENTER)

# ========== Slide 11: Temporal Alignment ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, 0.8, 0.4, 10, 0.8, "Token 特征的\"时代近视\"现象", 32, ACCENT2, bold=True)
add_accent_line(slide, 0.8, 1.1, 6)

# Timeline visualization
timeline = [
    ("2019-2020", "GPT-1/2/3", "Token AUC ≈ 0.50", ACCENT2),
    ("2022", "ChatGPT", "Token AUC ≈ 0.50", ACCENT2),
    ("2023", "Mistral-7B (观察者)", "—", GRAY),
    ("2023-2024", "GPT-4/Claude/Cohere", "Token AUC = 0.99", ACCENT3),
]
for i, (year, model, auc, clr) in enumerate(timeline):
    y = 1.8 + i * 1.2
    add_text_box(slide, 1.5, y, 2, 0.5, year, 20, clr, bold=True)
    add_text_box(slide, 4, y, 4, 0.5, model, 18, WHITE)
    add_text_box(slide, 8.5, y, 4, 0.5, auc, 18, clr, bold=True)

add_text_box(slide, 0.8, 6.2, 12, 0.8,
    "Mistral-7B 只能看到与自己同代的 AI \"指纹\" — 观察者模型需要及时更新",
    18, ACCENT, alignment=PP_ALIGN.CENTER)

# ========== Slide 12: Per-Generator ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, 0.8, 0.4, 10, 0.8, "Per-Generator 分析", 36, ACCENT, bold=True)
add_accent_line(slide, 0.8, 1.1, 4)
add_image_safe(slide, FIG / "raid_generator_radar.png", 0.2, 1.3, width=6.5)
add_image_safe(slide, FIG / "raid_120_per_generator.png", 6.8, 1.3, width=6.2)
add_text_box(slide, 0.8, 6.6, 11, 0.5,
    "不同生成器有不同的\"特征指纹\" — 120 维在 Cohere/ChatGPT 上提升最大",
    15, ACCENT, alignment=PP_ALIGN.CENTER)

# ========== Slide 13: Adversarial ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, 0.8, 0.4, 10, 0.8, "对抗攻击鲁棒性", 36, ACCENT, bold=True)
add_accent_line(slide, 0.8, 1.1, 4)
add_image_safe(slide, FIG / "raid_adversarial.png", 0.3, 1.3, width=7)

add_text_box(slide, 7.8, 1.5, 5, 0.6, "结论", 22, ACCENT, bold=True)
add_bullet_list(slide, 7.8, 2.2, 5, 4.5, [
    "10/11 种攻击几乎无效",
    "  统计特征对表面扰动有天然鲁棒性",
    "",
    "唯一有效攻击: Paraphrase",
    "  AUC 从 0.9951 降至 0.7902",
    "  因为释义改变了文本的宏观结构",
    "",
    "启示: 需要引入语义相似度特征",
    "  来防御释义攻击",
], font_size=16, color=LIGHT)

# ========== Slide 14: Cross-Dataset ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, 0.8, 0.4, 10, 0.8, "跨数据集验证 (5 数据集完整结果)", 32, ACCENT, bold=True)
add_accent_line(slide, 0.8, 1.1, 8)

# Full table
header_items = [("Dataset", 0.8, 2.5), ("90-feat", 3.5, 1.5), ("30-tok", 5.2, 1.5),
                ("120-mix", 6.9, 1.5), ("Delta", 8.6, 1.5), ("解读", 10.3, 2.5)]
for name, x, w in header_items:
    add_text_box(slide, x, 1.4, w, 0.5, name, 16, ACCENT, bold=True)

rows = [
    ("RAID (多代际)", "0.9951", "0.9900", "0.9992", "+0.004", "两侧有信号 → 互补", ACCENT3),
    ("SemEval (纯新AI)", "0.5136", "0.9766", "0.8443", "-0.132", "噪声稀释!", ACCENT2),
    ("TuringBench (纯老AI)", "0.9899", "0.4935", "0.9847", "-0.005", "token噪声少,影响小", LIGHT),
    ("HC3 (单模型)", "0.9996", "0.4986", "0.9996", "+0.000", "已饱和,无提升空间", GRAY),
    ("Pile (老模型)", "0.9735", "0.5026", "0.9727", "-0.001", "同上", GRAY),
]
for i, (ds, a90, a30, a120, delta, note, clr) in enumerate(rows):
    y = 2.1 + i * 0.75
    add_text_box(slide, 0.8, y, 2.5, 0.5, ds, 14, clr, bold=True)
    add_text_box(slide, 3.5, y, 1.5, 0.5, a90, 14, ACCENT2 if a90.startswith("0.51") else WHITE)
    add_text_box(slide, 5.2, y, 1.5, 0.5, a30, 14, ACCENT2 if a30.startswith("0.49") else WHITE)
    add_text_box(slide, 6.9, y, 1.5, 0.5, a120, 14, clr)
    add_text_box(slide, 8.6, y, 1.5, 0.5, delta, 14, ACCENT3 if delta.startswith("+0.004") else ACCENT2 if delta.startswith("-0.1") else LIGHT)
    add_text_box(slide, 10.3, y, 2.5, 0.5, note, 12, GRAY)

add_text_box(slide, 0.8, 6.0, 12, 0.8,
    "核心规律: 120 组合只在两侧特征都有信号时互补增强 (RAID)\n"
    "当一侧是纯噪声时，噪声特征越多稀释越严重 (SemEval: 88噪声 vs 30有效 → AUC -0.13)",
    15, ACCENT, alignment=PP_ALIGN.CENTER)

# ========== Slide 14b: Noise Dilution Deep Dive ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, 0.8, 0.4, 12, 0.8, "SemEval: 噪声特征稀释现象剖析", 32, ACCENT2, bold=True)
add_accent_line(slide, 0.8, 1.1, 8)

add_text_box(slide, 0.8, 1.4, 5.5, 0.6, "为什么 90-feat 在 SemEval 上失效?", 20, WHITE, bold=True)
add_bullet_list(slide, 0.8, 2.0, 5.5, 2.5, [
    "88 个特征中 没有一个 AUC > 0.55",
    "最好的 emb_pca_39 也只有 AUC = 0.5159",
    "",
    "说明 GPT-4/Claude 3 等新 AI 已经在",
    "词汇多样性、句法结构、可读性等所有",
    "表面统计维度上与人类文本无法区分",
    "",
    "→ 手工特征的\"保质期\"正在到来",
], font_size=15, color=LIGHT)

add_text_box(slide, 7, 1.4, 5.5, 0.6, "XGBoost 的过拟合机制", 20, WHITE, bold=True)
add_bullet_list(slide, 7, 2.0, 5.5, 2.5, [
    "120 维模型中 42.7% importance → 噪声特征",
    "  punct_semicolon_rate: imp=0.079 (纯噪声!)",
    "",
    "噪声特征在训练集上碰巧与标签相关",
    "XGBoost 把虚假相关当作信号",
    "测试集上虚假信号不成立 → 拖累预测",
    "",
    "→ 88 个噪声 vs 30 个有效 = 稀释严重",
], font_size=15, color=LIGHT)

# Bottom: progressive degradation
add_text_box(slide, 0.8, 5.0, 12, 0.5, "逐步加入 90-feat 的效果 (越多噪声越差):", 18, ACCENT, bold=True)
steps = [
    ("30 tok alone", "0.9766", ACCENT3),
    ("+ top5 90-feat (35维)", "0.9719", ACCENT3),
    ("+ top10 90-feat (40维)", "0.9436", RGBColor(0xFF, 0xC1, 0x07)),
    ("+ all 88 feat (118维)", "0.8443", ACCENT2),
]
for i, (label, auc, clr) in enumerate(steps):
    x = 0.8 + i * 3.1
    add_text_box(slide, x, 5.5, 2.8, 0.4, label, 14, LIGHT)
    add_text_box(slide, x, 5.9, 2.8, 0.5, f"AUC = {auc}", 20, clr, bold=True)

add_text_box(slide, 0.8, 6.6, 12, 0.5,
    "启示: 生产系统需要特征选择/门控 — 先评估各特征组信号强度，动态决定使用哪些特征",
    15, ACCENT, alignment=PP_ALIGN.CENTER)

# ========== Slide 14c: Auto-Filter Solution ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, 0.8, 0.4, 12, 0.8, "解决方案: Auto-Filter 噪声特征自动过滤", 32, ACCENT3, bold=True)
add_accent_line(slide, 0.8, 1.1, 8)

add_text_box(slide, 0.8, 1.4, 5.5, 0.6, "机制: 组级信号强度检测", 20, WHITE, bold=True)
add_bullet_list(slide, 0.8, 2.0, 5.5, 2.5, [
    "对每组特征计算平均单特征 AUC",
    "  avg_AUC = mean(max(AUC_i, 1-AUC_i))",
    "丢弃 avg AUC < 0.52 的组 (近随机)",
    "余下特征 + 30 token → 组合训练",
    "",
    "只用训练集评估 → 不泄露测试信息",
    "O(n·p) 开销 → 几秒完成",
], font_size=15, color=LIGHT)

add_text_box(slide, 7, 1.4, 5.5, 0.6, "跨数据集 Auto-Filter 结果", 20, WHITE, bold=True)

# Table
af_header = [("Dataset", 7.0, 2.2), ("Before", 9.3, 1.3), ("After", 10.7, 1.3)]
for name, x, w in af_header:
    add_text_box(slide, x, 2.0, w, 0.4, name, 15, ACCENT, bold=True)

af_rows = [
    ("RAID", "0.9992", "0.9992", "drop 2 组 (无影响)", LIGHT),
    ("SemEval", "0.8443", "0.9763", "drop 全部9组 → tok-only", ACCENT3),
    ("TuringBench", "0.9847", "0.9847", "drop 0 组 (无影响)", LIGHT),
]
for i, (ds, before, after, note, clr) in enumerate(af_rows):
    y = 2.6 + i * 0.65
    add_text_box(slide, 7.0, y, 2.2, 0.4, ds, 14, clr)
    add_text_box(slide, 9.3, y, 1.3, 0.4, before, 14, WHITE)
    add_text_box(slide, 10.7, y, 1.3, 0.4, after, 14, ACCENT3 if after != before else WHITE)

add_text_box(slide, 7, 4.8, 5.5, 0.4, note if False else "", 12, GRAY)

# Notes for each row
for i, (_, _, _, note, _) in enumerate(af_rows):
    y = 2.6 + i * 0.65
    add_text_box(slide, 7.0, y + 0.3, 5.5, 0.3, note, 11, GRAY)

add_text_box(slide, 0.8, 5.5, 12, 0.8,
    "Auto-Filter 是\"只赚不赔\"策略: RAID/TuringBench 性能不变, SemEval AUC 从 0.84 恢复到 0.98\n"
    "已集成到项目 pipeline: src/auto_filter.py",
    16, ACCENT, alignment=PP_ALIGN.CENTER)

# ========== Slide 15: RAID vs HC3 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, 0.8, 0.4, 10, 0.8, "RAID vs HC3: 特征组消融对比", 32, ACCENT, bold=True)
add_accent_line(slide, 0.8, 1.1, 6)

# Table
header = [("特征组", 1), ("RAID AUC", 4.5), ("HC3 AUC", 7), ("启示", 9.5)]
for name, x in header:
    add_text_box(slide, x, 1.5, 2.5, 0.5, name, 18, ACCENT, bold=True)

data = [
    ("basic_counts", "0.9700", "0.9204", "RAID 更强"),
    ("averages", "0.9425", "0.9021", "RAID 更强"),
    ("lexical_richness", "0.8861", "0.9367", "HC3 更强"),
    ("readability", "0.8500", "0.9741", "HC3 远优"),
    ("perplexity", "0.4920", "0.9912", "完全反转!"),
    ("embedding_pca", "0.9357", "0.8972", "接近"),
]
for i, (g, r_auc, h_auc, note) in enumerate(data):
    y = 2.2 + i * 0.7
    clr = ACCENT2 if "反转" in note else WHITE
    add_text_box(slide, 1, y, 3, 0.5, g, 16, clr)
    add_text_box(slide, 4.5, y, 2, 0.5, r_auc, 16, clr)
    add_text_box(slide, 7, y, 2, 0.5, h_auc, 16, clr)
    add_text_box(slide, 9.5, y, 3, 0.5, note, 16, clr)

add_text_box(slide, 0.8, 6.3, 11.5, 0.8,
    "核心启示: 多模型检测不能依赖单一语言模型的困惑度 — 需要模型无关的统计特征作为稳定基线",
    16, ACCENT, alignment=PP_ALIGN.CENTER)

# ========== Slide 16: Conclusion ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, 0.8, 0.4, 10, 0.8, "结论", 36, ACCENT, bold=True)
add_accent_line(slide, 0.8, 1.1, 4)

conclusions = [
    "1. XGBoost + 120 维组合特征在 RAID 上达到 AUC 0.9992",
    "    完全可解释，超越零样本方法 Fast-DetectGPT (0.7815)",
    "",
    "2. 两类特征视角正交，互补覆盖盲区",
    "    手工特征看\"文本外表\"，Token 特征看\"AI 内在\"",
    "",
    "3. GPT-2 困惑度的 Paradox 是重要发现",
    "    HC3 最强 (0.99) → RAID 完全失效 (0.49)",
    "",
    "4. 统计特征对 10/11 种对抗攻击天然鲁棒",
    "    只有释义攻击能显著降低检测效果",
    "",
    "5. Token 概率特征有\"时代窗口\"",
    "    观察者模型需要与目标生成器同代",
]
add_bullet_list(slide, 0.8, 1.5, 11, 5.5, conclusions, font_size=18, color=LIGHT, spacing=0.8)

# ========== Slide 17: Future Work ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, 0.8, 0.4, 10, 0.8, "未来方向与部署建议", 36, ACCENT, bold=True)
add_accent_line(slide, 0.8, 1.1, 4)

add_text_box(slide, 0.8, 1.5, 5.5, 0.6, "未来工作", 22, ACCENT3, bold=True)
add_bullet_list(slide, 0.8, 2.2, 5.5, 3, [
    "引入语义相似度特征防御释义攻击",
    "多观察者模型集成 (Mistral + Qwen + Llama)",
    "动态特征选择: 根据文本领域自动调整",
    "探索更多 token 统计: burstiness, curvature",
    "在 SemEval/TuringBench 全量数据验证",
], font_size=16, color=LIGHT)

add_text_box(slide, 7, 1.5, 5.5, 0.6, "部署建议", 22, ACCENT2, bold=True)
add_bullet_list(slide, 7, 2.2, 5.5, 3, [
    "以模型无关的手工特征作为稳定基线",
    "辅以同代 LLM token 特征覆盖新生成器",
    "观察者模型需要定期更新 (半年一次)",
    "跨域部署需注意领域差异 (诗歌 vs 新闻)",
    "提供 SHAP 解释以增强用户信任",
], font_size=16, color=LIGHT)

add_text_box(slide, 0.8, 5.5, 11.5, 1.2,
    "\"手工特征看文本的外表，Token 概率特征看文本在 AI 眼中的内在。\n"
    "当一个视角被欺骗时，另一个视角往往能看穿。这就是 120 维组合的力量。\"",
    20, ACCENT, alignment=PP_ALIGN.CENTER)

# ========== Slide 18: Thank You ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, 1, 2.5, 11, 1.5,
    "Thank You", 52, ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
add_accent_line(slide, 4, 4.0, 5)
add_text_box(slide, 1, 4.3, 11, 1,
    "Q & A", 32, WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1, 5.5, 11, 1,
    "代码: github.com/HaibinLai/HC3  |  赖海斌  |  CS306 数据挖掘", 16, GRAY, alignment=PP_ALIGN.CENTER)

# Save
out_path = OUT / "AIGC_Detection_Presentation.pptx"
prs.save(str(out_path))
print(f"Saved: {out_path}")
print(f"Slides: {len(prs.slides)}")
